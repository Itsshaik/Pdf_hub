from django.shortcuts import render, redirect
from django.http import HttpResponse, Http404, JsonResponse
from django.contrib import messages
from django.core.files.base import ContentFile
from django.core.files.storage import default_storage
from .models import ProcessedFile
from PIL import Image
import PyPDF2
from pptx import Presentation
from pptx.util import Inches
import pytesseract
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from pdf2image import convert_from_path
import io
import os
import tempfile

def home(request):
    """Main page with all conversion options"""
    return render(request, 'converter/home.html')

def image_to_pdf_view(request):
    """Convert images to PDF"""
    if request.method == 'POST':
        uploaded_files = request.FILES.getlist('images')
        
        if not uploaded_files:
            messages.error(request, 'Please select at least one image file.')
            return redirect('image_to_pdf')
        
        try:
            # Create PDF from images
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            output = io.BytesIO()
            c = canvas.Canvas(output, pagesize=letter)
            
            for uploaded_file in uploaded_files:
                # Process each image
                image = Image.open(uploaded_file)
                
                # Convert to RGB if necessary
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Save image temporarily
                temp_image = io.BytesIO()
                image.save(temp_image, format='JPEG')
                temp_image.seek(0)
                
                # Calculate dimensions to fit on page
                img_width, img_height = image.size
                page_width, page_height = letter
                
                # Scale image to fit page while maintaining aspect ratio
                scale = min(page_width / img_width, page_height / img_height)
                scaled_width = img_width * scale
                scaled_height = img_height * scale
                
                # Center image on page
                x = (page_width - scaled_width) / 2
                y = (page_height - scaled_height) / 2
                
                # Create temporary file for image
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                    image.save(temp_file.name, 'JPEG')
                    c.drawImage(temp_file.name, x, y, scaled_width, scaled_height)
                    os.unlink(temp_file.name)
                
                c.showPage()
            
            c.save()
            output.seek(0)
            
            # Save to model
            processed_file = ProcessedFile.objects.create(
                operation='image_to_pdf',
                original_file=uploaded_files[0],  # Store first file as reference
                file_size=sum(f.size for f in uploaded_files),
                is_processed=True
            )
            
            # Save PDF
            pdf_content = ContentFile(output.getvalue())
            processed_file.processed_file.save(
                f'converted_{processed_file.id}.pdf',
                pdf_content
            )
            
            messages.success(request, 'Images successfully converted to PDF!')
            return redirect('download_file', file_id=processed_file.id)
            
        except Exception as e:
            messages.error(request, f'Error converting images: {str(e)}')
            return redirect('image_to_pdf')
    
    return render(request, 'converter/image_to_pdf.html')

def pdf_merge_view(request):
    """Merge multiple PDFs into one"""
    if request.method == 'POST':
        uploaded_files = request.FILES.getlist('pdfs')
        
        if len(uploaded_files) < 2:
            messages.error(request, 'Please select at least two PDF files to merge.')
            return redirect('pdf_merge')
        
        try:
            merger = PyPDF2.PdfMerger()
            
            for uploaded_file in uploaded_files:
                # Verify it's a PDF
                if not uploaded_file.name.lower().endswith('.pdf'):
                    messages.error(request, f'{uploaded_file.name} is not a PDF file.')
                    return redirect('pdf_merge')
                
                merger.append(uploaded_file)
            
            output = io.BytesIO()
            merger.write(output)
            merger.close()
            output.seek(0)
            
            # Save to model
            processed_file = ProcessedFile.objects.create(
                operation='pdf_merge',
                original_file=uploaded_files[0],  # Store first file as reference
                file_size=sum(f.size for f in uploaded_files),
                is_processed=True
            )
            
            # Save merged PDF
            pdf_content = ContentFile(output.getvalue())
            processed_file.processed_file.save(
                f'merged_{processed_file.id}.pdf',
                pdf_content
            )
            
            messages.success(request, 'PDFs successfully merged!')
            return redirect('download_file', file_id=processed_file.id)
            
        except Exception as e:
            messages.error(request, f'Error merging PDFs: {str(e)}')
            return redirect('pdf_merge')
    
    return render(request, 'converter/pdf_merge.html')

def pdf_password_view(request):
    """Add password protection to PDF"""
    if request.method == 'POST':
        uploaded_file = request.FILES.get('pdf')
        password = request.POST.get('password')
        
        if not uploaded_file or not password:
            messages.error(request, 'Please select a PDF file and enter a password.')
            return redirect('pdf_password')
        
        try:
            # Read the original PDF
            reader = PyPDF2.PdfReader(uploaded_file)
            writer = PyPDF2.PdfWriter()
            
            # Copy all pages
            for page in reader.pages:
                writer.add_page(page)
            
            # Encrypt with password
            writer.encrypt(password)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            # Save to model
            processed_file = ProcessedFile.objects.create(
                operation='pdf_password',
                original_file=uploaded_file,
                file_size=uploaded_file.size,
                is_processed=True
            )
            
            # Save encrypted PDF
            pdf_content = ContentFile(output.getvalue())
            processed_file.processed_file.save(
                f'encrypted_{processed_file.id}.pdf',
                pdf_content
            )
            
            messages.success(request, 'PDF successfully password protected!')
            return redirect('download_file', file_id=processed_file.id)
            
        except Exception as e:
            messages.error(request, f'Error protecting PDF: {str(e)}')
            return redirect('pdf_password')
    
    return render(request, 'converter/pdf_password.html')

def pdf_to_ppt_view(request):
    """Convert PDF to PowerPoint"""
    if request.method == 'POST':
        uploaded_file = request.FILES.get('pdf')
        
        if not uploaded_file:
            messages.error(request, 'Please select a PDF file.')
            return redirect('pdf_to_ppt')
        
        try:
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                for chunk in uploaded_file.chunks():
                    temp_pdf.write(chunk)
                temp_pdf_path = temp_pdf.name
            
            # Convert PDF pages to images
            images = convert_from_path(temp_pdf_path)
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            for i, image in enumerate(images):
                # Add slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Save image temporarily
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                    image.save(temp_img.name, 'PNG')
                    
                    # Add image to slide
                    slide.shapes.add_picture(temp_img.name, Inches(0), Inches(0), 
                                           width=Inches(10), height=Inches(7.5))
                    
                    os.unlink(temp_img.name)
            
            # Save PowerPoint
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            # Clean up temp PDF
            os.unlink(temp_pdf_path)
            
            # Save to model
            processed_file = ProcessedFile.objects.create(
                operation='pdf_to_ppt',
                original_file=uploaded_file,
                file_size=uploaded_file.size,
                is_processed=True
            )
            
            # Save PowerPoint
            ppt_content = ContentFile(output.getvalue())
            processed_file.processed_file.save(
                f'presentation_{processed_file.id}.pptx',
                ppt_content
            )
            
            messages.success(request, 'PDF successfully converted to PowerPoint!')
            return redirect('download_file', file_id=processed_file.id)
            
        except Exception as e:
            messages.error(request, f'Error converting PDF to PowerPoint: {str(e)}')
            return redirect('pdf_to_ppt')
    
    return render(request, 'converter/pdf_to_ppt.html')

def ocr_extraction_view(request):
    """Extract text from images using OCR"""
    if request.method == 'POST':
        uploaded_file = request.FILES.get('image')
        
        if not uploaded_file:
            messages.error(request, 'Please select an image file.')
            return redirect('ocr_extraction')
        
        try:
            # Open image
            image = Image.open(uploaded_file)
            
            # Extract text using OCR
            extracted_text = pytesseract.image_to_string(image)
            
            if not extracted_text.strip():
                extracted_text = "No text found in the image."
            
            # Save to model
            processed_file = ProcessedFile.objects.create(
                operation='ocr_extraction',
                original_file=uploaded_file,
                extracted_text=extracted_text,
                file_size=uploaded_file.size,
                is_processed=True
            )
            
            messages.success(request, 'Text successfully extracted from image!')
            return render(request, 'converter/ocr_result.html', {
                'extracted_text': extracted_text,
                'file_id': processed_file.id
            })
            
        except Exception as e:
            messages.error(request, f'Error extracting text: {str(e)}')
            return redirect('ocr_extraction')
    
    return render(request, 'converter/ocr_extraction.html')

def download_file(request, file_id):
    """Download processed file"""
    try:
        processed_file = ProcessedFile.objects.get(id=file_id)
        
        if not processed_file.processed_file:
            raise Http404("Processed file not found.")
        
        file_path = processed_file.processed_file.path
        
        if not os.path.exists(file_path):
            raise Http404("File not found on disk.")
        
        with open(file_path, 'rb') as f:
            response = HttpResponse(f.read())
            
        # Set appropriate content type
        if file_path.endswith('.pdf'):
            content_type = 'application/pdf'
        elif file_path.endswith('.pptx'):
            content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            content_type = 'application/octet-stream'
        
        response['Content-Type'] = content_type
        response['Content-Disposition'] = f'attachment; filename="{os.path.basename(file_path)}"'
        
        return response
        
    except ProcessedFile.DoesNotExist:
        raise Http404("File not found.")

def file_history(request):
    """Show processing history"""
    files = ProcessedFile.objects.all().order_by('-created_at')[:20]
    return render(request, 'converter/history.html', {'files': files})
